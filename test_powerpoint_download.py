#!/usr/bin/env python3
"""
Test script to verify PowerPoint creation and download functionality
"""

import requests
import json
import time

# Debug server URL
DEBUG_URL = "http://localhost:8069"

def test_powerpoint_creation():
    """Test creating a PowerPoint through the debug interface."""
    print("🧪 Testing PowerPoint creation and download flow...")

    # Test data
    conversation_id = "test-ppt-session"

    # 1. Test ZSH command to create PowerPoint
    print("\n1. Creating PowerPoint via ZSH container...")
    ppt_command = """python3 -c "
from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Test PowerPoint'
slide.placeholders[1].text = 'Created via Debug Interface\\nWith automatic download links!'
prs.save('/mnt/test_presentation.pptx')
print('PowerPoint created successfully!')
" """

    tool_data = {
        "tool_id": "zsh",
        "parameters": {"command": ppt_command},
        "conversation_id": conversation_id
    }

    try:
        response = requests.post(f"{DEBUG_URL}/api/tool/execute", json=tool_data)
        response.raise_for_status()
        result = response.json()

        print(f"✅ Tool execution status: {result.get('error', 'success')}")

        if 'downloadable_files' in result.get('result', {}):
            files = result['result']['downloadable_files']
            print(f"📁 Found {files.get('total_files', 0)} downloadable files")

            if 'powerpoint_files' in files:
                ppt_files = files['powerpoint_files']
                print(f"🎯 PowerPoint files: {len(ppt_files)}")
                for file in ppt_files:
                    print(f"   📄 {file['name']} - {file['download_url']}")

                    # Test download
                    print(f"\n2. Testing download of {file['name']}...")
                    download_response = requests.get(f"{DEBUG_URL}{file['download_url']}")
                    if download_response.status_code == 200:
                        print(f"✅ Download successful! Size: {len(download_response.content)} bytes")
                    else:
                        print(f"❌ Download failed: {download_response.status_code}")

        # 3. Test file listing API
        print(f"\n3. Testing file listing API...")
        files_response = requests.get(f"{DEBUG_URL}/api/files/container/{conversation_id}")
        if files_response.status_code == 200:
            files_data = files_response.json()
            print(f"✅ File listing API works! Found {len(files_data.get('files', []))} files")
            for file in files_data.get('files', []):
                if file.get('extension') == '.pptx':
                    print(f"   🎯 {file['name']} ({file['size']} bytes)")
        else:
            print(f"❌ File listing failed: {files_response.status_code}")

    except Exception as e:
        print(f"❌ Test failed: {e}")
        return False

    print("\n🎉 PowerPoint download flow test completed!")
    return True

if __name__ == "__main__":
    test_powerpoint_creation()
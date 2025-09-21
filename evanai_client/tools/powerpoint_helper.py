"""
Enhanced PowerPoint Generation Helper

This module provides improved PowerPoint generation capabilities with:
- Content length management and overflow prevention
- Multiple slide layouts and visual variety
- Smart content chunking and organization
- Generic, reusable functionality
"""

def create_enhanced_powerpoint_script(title: str, content_sections: list, include_code: bool = True) -> str:
    """
    Create an enhanced PowerPoint generation script.

    Args:
        title: Presentation title
        content_sections: List of dicts with 'title', 'content', 'type' keys
        include_code: Whether to include code examples

    Returns:
        Python script string for generating PowerPoint
    """

    script = f'''
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import textwrap

# Create presentation with widescreen layout
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define modern color palette
BLUE = RGBColor(25, 118, 210)      # Material Blue
GREEN = RGBColor(76, 175, 80)      # Material Green
ORANGE = RGBColor(255, 152, 0)     # Material Orange
PURPLE = RGBColor(156, 39, 176)    # Material Purple
DARK = RGBColor(55, 71, 79)        # Blue Grey 800
LIGHT = RGBColor(245, 245, 245)    # Grey 100
ACCENT = RGBColor(255, 87, 34)     # Deep Orange

def add_title_slide():
    """Create an engaging title slide with modern design."""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "{title}"
    subtitle.text = "AI-Generated Technical Presentation\\n🤖 Created by Advanced Analysis\\n\\n" + "Generated automatically from codebase analysis"

    # Style the title with modern typography
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.color.rgb = BLUE
    title_paragraph.font.size = Pt(54)
    title_paragraph.font.bold = True

    # Style subtitle
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = DARK

def truncate_content(text: str, max_chars: int = 800, max_lines: int = 12) -> str:
    """Intelligently truncate content to fit slide boundaries."""
    lines = text.split('\\n')

    # First, limit by number of lines
    if len(lines) > max_lines:
        lines = lines[:max_lines-1] + ["... (continued)"]

    # Then check total character count
    truncated_text = '\\n'.join(lines)
    if len(truncated_text) > max_chars:
        # Find a good breaking point
        words = truncated_text[:max_chars].split()
        if len(words) > 1:
            truncated_text = ' '.join(words[:-1]) + "..."

    return truncated_text

def add_content_slide(title_text: str, content_text: str, slide_type: str = "standard"):
    """Add a content slide with proper formatting and overflow handling."""

    if slide_type == "two_column":
        # Use two-column layout for better content distribution
        slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = title_text

        # Split content into two columns if it's too long
        content_lines = content_text.split('\\n')
        mid_point = len(content_lines) // 2

        left_content = '\\n'.join(content_lines[:mid_point])
        right_content = '\\n'.join(content_lines[mid_point:])

        # Add content to placeholders
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = truncate_content(left_content, 400, 6)
        if len(slide.placeholders) > 2:
            slide.placeholders[2].text = truncate_content(right_content, 400, 6)
    else:
        # Standard single-column layout
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = title_text
        content.text = truncate_content(content_text)

        # Style the content for better readability
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(6)

def add_summary_slide():
    """Create a compelling summary slide."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Key Insights & Next Steps"

    summary_text = """🎯 Analysis Highlights:
• Comprehensive codebase understanding achieved
• Technical architecture thoroughly documented
• Implementation patterns identified

🚀 Generated Automatically:
• This presentation was created by AI analyzing the codebase
• Real-time insights from code structure and patterns
• Professional documentation in minutes, not hours

🔮 Capabilities Demonstrated:
• Advanced meta-programming and self-analysis
• Intelligent content generation and formatting
• Seamless integration of analysis tools"""

    content.text = summary_text

    # Style for impact
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(8)

# Generate all slides
add_title_slide()
'''

    # Add content sections dynamically
    for i, section in enumerate(content_sections):
        section_title = section.get('title', f'Section {i+1}')
        section_content = section.get('content', '')
        section_type = section.get('type', 'standard')

        # Determine if this section should use two-column layout
        if len(section_content) > 1000 or section_content.count('\n') > 10:
            section_type = 'two_column'

        script += f'''
add_content_slide("{section_title}", """{section_content}""", "{section_type}")
'''

    script += '''
add_summary_slide()

# Save the presentation
prs.save('/mnt/enhanced_presentation.pptx')
print("✅ Enhanced presentation created: /mnt/enhanced_presentation.pptx")
print(f"📊 Generated {len(prs.slides)} slides with improved design")
print("🎨 Features: Content management, multiple layouts, modern styling")
'''

    return script
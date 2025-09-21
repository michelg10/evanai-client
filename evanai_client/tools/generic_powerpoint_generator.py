#!/usr/bin/env python3
"""
Generic PowerPoint Generation Script

This script can be used in containers to generate professional PowerPoint presentations
from any analysis content. It features:
- Enhanced visual design with multiple layouts
- Content overflow management
- Smart content chunking
- Professional styling

Usage in container:
    python3 generic_powerpoint_generator.py --title "My Presentation" --analysis "content here"
"""

import argparse
import json
import sys
from pathlib import Path

def create_powerpoint_script(title: str, analysis_content: str, output_file: str = "/mnt/presentation.pptx") -> str:
    """Create enhanced PowerPoint generation script."""

    # Extract key sections from analysis content
    sections = parse_analysis_content(analysis_content)

    script = f'''
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import textwrap
import re

# Create presentation with widescreen layout
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define modern color palette
BLUE = RGBColor(25, 118, 210)      # Material Blue
GREEN = RGBColor(76, 175, 80)      # Material Green
ORANGE = RGBColor(255, 152, 0)     # Material Orange
PURPLE = RGBColor(156, 39, 176)    # Material Purple
DARK = RGBColor(55, 71, 79)        # Blue Grey 800
LIGHT = RGBColor(245, 245, 245)    # Grey 100

def truncate_content(text: str, max_chars: int = 700, max_lines: int = 10) -> str:
    """Intelligently truncate content to fit slide boundaries."""
    if not text:
        return ""

    lines = text.split('\\n')
    # Remove empty lines and clean up
    lines = [line.strip() for line in lines if line.strip()]

    # Limit by number of lines first
    if len(lines) > max_lines:
        lines = lines[:max_lines-1] + ["... (content truncated for readability)"]

    # Then check total character count
    truncated_text = '\\n'.join(lines)
    if len(truncated_text) > max_chars:
        # Find a good breaking point at sentence or bullet point
        truncated = truncated_text[:max_chars]
        if '\\n' in truncated:
            lines = truncated.split('\\n')
            truncated_text = '\\n'.join(lines[:-1]) + '\\n... (continued)'
        else:
            # Break at last complete sentence
            sentences = truncated.split('. ')
            if len(sentences) > 1:
                truncated_text = '. '.join(sentences[:-1]) + '...'
            else:
                truncated_text = truncated + '...'

    return truncated_text

def add_title_slide():
    """Create an engaging title slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "{title}"
    subtitle.text = "AI-Generated Technical Analysis\\n🤖 Automated Documentation\\n\\nCreated from comprehensive codebase analysis"

    # Style the title
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.color.rgb = BLUE
    title_paragraph.font.size = Pt(48)
    title_paragraph.font.bold = True

def add_content_slide(title_text: str, content_text: str, use_two_columns: bool = False):
    """Add a content slide with smart formatting."""

    if use_two_columns and len(content_text) > 600:
        # Try to use two-column layout if available
        slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    else:
        slide_layout = prs.slide_layouts[1]  # Title and content

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

    # Style title
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.color.rgb = DARK
    title_paragraph.font.bold = True

    if use_two_columns and len(slide.placeholders) > 2:
        # Split content for two columns
        content_lines = content_text.split('\\n')
        mid_point = len(content_lines) // 2

        left_content = '\\n'.join(content_lines[:mid_point])
        right_content = '\\n'.join(content_lines[mid_point:])

        slide.placeholders[1].text = truncate_content(left_content, 350, 6)
        slide.placeholders[2].text = truncate_content(right_content, 350, 6)
    else:
        # Single column
        content = slide.placeholders[1]
        content.text = truncate_content(content_text)

        # Style content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(6)

def add_summary_slide():
    """Create a professional summary slide."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Analysis Summary"

    summary_text = """🎯 Key Findings:
• Comprehensive technical analysis completed
• Architecture and implementation patterns documented
• Professional insights generated automatically

🤖 Generated by AI:
• This presentation was created through automated code analysis
• Real-time insights from project structure and implementation
• Professional documentation generated in seconds

✨ Demonstrates:
• Advanced AI capabilities for technical documentation
• Intelligent content analysis and presentation generation
• Seamless integration of analysis and visualization tools"""

    content.text = summary_text

# Generate all slides
add_title_slide()
'''

    # Add sections dynamically
    for section in sections:
        title = section['title']
        content = section['content']
        use_two_cols = len(content) > 800  # Use two columns for lengthy content

        script += f'''
add_content_slide("{title}", """{content}""", {use_two_cols})
'''

    script += f'''
add_summary_slide()

# Save the presentation
prs.save('{output_file}')
print("✅ Enhanced presentation created: {output_file}")
print(f"📊 Generated {{len(prs.slides)}} slides with professional design")
print("🎨 Features: Smart content management, multiple layouts, modern styling")
'''

    return script

def parse_analysis_content(content: str) -> list:
    """Parse analysis content into structured sections for slides."""
    sections = []

    # Try to detect markdown-style headers
    if '##' in content or '#' in content:
        current_section = None
        current_content = []

        for line in content.split('\n'):
            line = line.strip()
            if line.startswith('##'):
                # Save previous section
                if current_section:
                    sections.append({
                        'title': current_section,
                        'content': '\\n'.join(current_content)
                    })

                # Start new section
                current_section = line.replace('##', '').replace('#', '').strip()
                current_content = []
            elif line.startswith('#') and not line.startswith('##'):
                # Main header - could be presentation title or major section
                if not current_section:
                    current_section = line.replace('#', '').strip()
                    current_content = []
            else:
                if line:  # Skip empty lines
                    current_content.append(line)

        # Add final section
        if current_section and current_content:
            sections.append({
                'title': current_section,
                'content': '\\n'.join(current_content)
            })

    # If no markdown headers found, split by paragraphs or create generic sections
    if not sections:
        # Split content into chunks
        content_chunks = content.split('\\n\\n')
        for i, chunk in enumerate(content_chunks[:6]):  # Limit to 6 sections
            if chunk.strip():
                sections.append({
                    'title': f'Analysis Section {i+1}',
                    'content': chunk.strip()
                })

    # Ensure we have at least one section
    if not sections:
        sections.append({
            'title': 'Technical Analysis',
            'content': content[:1000] + '...' if len(content) > 1000 else content
        })

    return sections

def main():
    parser = argparse.ArgumentParser(description='Generate enhanced PowerPoint presentations')
    parser.add_argument('--title', required=True, help='Presentation title')
    parser.add_argument('--analysis', required=True, help='Analysis content for slides')
    parser.add_argument('--output', default='/mnt/presentation.pptx', help='Output file path')

    args = parser.parse_args()

    print(f"🎨 Generating enhanced PowerPoint: {{args.title}}")

    # Create the PowerPoint script
    script = create_powerpoint_script(args.title, args.analysis, args.output)

    # Execute the script
    exec(script)

if __name__ == '__main__':
    main()